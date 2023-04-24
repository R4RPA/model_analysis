import math

from pptx.shapes.graphfrm import GraphicFrame
import os, shutil
import glob
import copy
from datetime import datetime
from copy import deepcopy
import collections
import collections.abc
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.table import Table, _Row, _Column, _Cell
from pptx.dml.color import ColorFormat, RGBColor
from time import sleep
import ctypes
from PIL import ImageGrab
import re
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

def get_bullet_info(paragraph, run=None):
    """Get info about bullet text format to reapply to another part of text"""
    pPr = paragraph._p.get_or_add_pPr()
    if run is None:
        run = paragraph.runs[0]
    p_info = {"marL": pPr.attrib['marL'], "indent": pPr.attrib['indent'], "level": paragraph.level,
              "fontName": run.font.name, "fontSize": run.font.size, }
    return p_info


def _sub_element(parent, tagname, **kwargs):
    """Helper for Paragraph bullet Point"""
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def set_bullet_paragraph(paragraph, bullets_info):
    """Apply bullets text format for given paragraph"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('marL', bullets_info['marL'])
    pPr.set('indent', bullets_info['indent'])
    _ = _sub_element(parent=pPr, tagname="a:buSzPct", val="100000")
    _ = _sub_element(parent=pPr, tagname="a:buFont", typeface=bullets_info['fontName'])
    _ = _sub_element(parent=pPr, tagname='a:buChar', char="•")


def move_slide(prs, old_index, new_index):
    """Move or re-arrange the slide to desired position"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def check_slide_layouts(prs):
    dest_slide_index = 0
    for num in range(100):
        for num2 in range(100):
            try:
                nc_layout = prs.slide_masters[num].slide_layouts[num2]
                copy_slide_index = 0
                dest_slide_index += 1
                prs = duplicate_slide(prs, nc_layout, copy_slide_index, dest_slide_index, '')
                print(num, num2, dest_slide_index)
            except:
                pass


def duplicate_slide(prs, layout, copy_slide_index, dest_slide_index, rst_type='', nc_list={}):
    """Duplicate the slide with the given index in presentation.
    Adds slide to given destination index position of the presentation"""

    """Original slide"""
    source = prs.slides[copy_slide_index]
    """Add a blank slide with selected layout"""
    dest = prs.slides.add_slide(layout)
    slide_id = prs.slides.index(dest)
    """Move or re-arrange the slide to desired position"""
    move_slide(prs, slide_id, dest_slide_index)

    """delete any default shapes are in the slide"""
    for shape in dest.shapes:
        sp = shape._sp
        sp.getparent().remove(sp)
    
    SkipThiShape = ''
    """Copy each shape in the original slide to duplicated slide"""

    if nc_list:
        for index, nc_dict in enumerate(nc_list):
            for key in nc_dict:
                for shape in source.shapes:
                    key2 = key.replace('{', '').replace('}', '')
                    if 'image' == key2.lower():
                        if shape.name == f'{rst_type}_image_{index+1}':
                            dest.shapes.add_picture(nc_dict[key], shape.left, shape.top, shape.width, shape.height)
                            SkipThiShape += shape.name
                            break

    for shape in source.shapes:
        if shape.name not in SkipThiShape:
            if 'Picture' in shape.name:
                """if shape has image, save the image to local, add to slide in position and delete"""
                dummy_image = shape.name + '.jpg'
                with open(dummy_image, 'wb') as f:
                    f.write(shape.image.blob)
                dest.shapes.add_picture(dummy_image, shape.left, shape.top, shape.width, shape.height)
                os.remove(dummy_image)
            elif f'{rst_type}_image' not in shape.name:
                if len(nc_list) < 6:
                    num_missing_elements = 6 - len(nc_list)
                    nc_list += [{}] * num_missing_elements
                if isinstance(shape, Table) or isinstance(shape, GraphicFrame):
                    start_index = 0
                    col_index = 0
                    number_of_slices = math.ceil(len(nc_list)/3)
                    for slice_num in range(number_of_slices):
                        slice_ = nc_list[start_index:start_index + 3]
                        row_index = 1
                        for index, nc_dict in enumerate(slice_):
                            update_table_row_2(table=shape if isinstance(shape, Table) else shape.table,
                                               nc_result=nc_dict, row_num=row_index, col_num=col_index)
                            row_index += 2
                        start_index += 3
                        col_index += 1

                newel = copy.deepcopy(shape.element)
                dest.shapes._spTree.insert_element_before(newel, 'p:extLst')

    #update_table_row_2()

    return prs


def update_table_row_2(table: Table, nc_result: dict, row_num: int = 0, font_size=12, col_num=0):
    """Update table row with nc results by tag identified"""
    cell = table.cell(row_num, col_num)
    cell.text = 'Mode ' + str(nc_result['mode']) + ', ' + str(nc_result['frequency']) + " Hz" if len(nc_result) > 0 else ''
    cell.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
    cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    cell.text_frame.word_wrap = False
    return table


def delete_slide(prs, delete_slide_index: list):
    """Delete a slide from presentation"""
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    for index in delete_slide_index:
        xml_slides.remove(slides[index])
    
    return prs


def get_object_by_name(slide, object_name):
    """Search and return the shape by name"""
    for shape in slide.shapes:
        if shape.name == object_name:
            if 'Table' in object_name:
                return shape.table
            else:
                return shape
    return None


def add_multiple_rows(table: Table, number_of_rows: int = 1, add_empty_row: bool = False):
    """Add multiple rows into table object"""
    for row in range(number_of_rows):
        table = add_row(table, add_empty_row)
    return table


def add_row(table: Table, add_empty_row: bool = False):
    """Add single row into table object"""
    new_row = deepcopy(table._tbl.tr_lst[-1])
    if add_empty_row:
        """if add_empty_row == True, then clear the cell content"""
        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            cell.text = ''
    table._tbl.append(new_row)
    return table


def remove_row(table: Table, row_to_delete: _Row):
    """Delete Table Row"""
    table._tbl.remove(row_to_delete._tr)


def replace_text_in_slide_tables(slide, nc_result: dict):
    """Replace tags in slide with nc results"""
    for shape in slide.shapes:
        try:
            table = shape.table
            update_table_values(table, nc_result)
        except:
            pass


def update_table_values(table: Table, nc_result: dict):
    row_count = len(table.rows)
    for row in range(-1, row_count-1):
        update_table_row(table, nc_result, row, font_size=12)


def update_table_row(table: Table, nc_result: dict, row_num: int = 0, font_size=12):
    """Update table row with nc results by tag identified"""
    cell1 = table.cell(row_num + 1, 0)
    cell1.text = str(nc_result['mode'])
    cell1.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    cell2 = table.cell(row_num + 1, 1)
    cell2.text = str(nc_result['frequency'])
    cell2.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    cells = [cell1, cell2]
    for cell in cells:
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = Pt(font_size)
        cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        cell.text_frame.word_wrap = False
    return table




def replace_text_in_prs(prs, nc_result: dict):
    for slide in prs.slides:
        replace_text_in_slide(slide, nc_result)


def replace_text_in_slide(slide, nc_result: dict):
    """Replace tags in slide with nc results"""
    green_txt = 'margin is in acceptable limits. Hence, this NC can be accepted from clearance standpoint.'
    red_txt = 'margin is not in acceptable limits. Hence, additional study is required.'
    for shape in slide.shapes:
        replace_text_in_shape(shape, nc_result)
        if shape.has_text_frame:
            if shape.text and green_txt in shape.text:
                shape.fill.solid()
                shape.fill.fore_color.rgb  = RGBColor(200, 230, 170)
            elif shape.text and red_txt in shape.text:
                shape.fill.solid()
                shape.fill.fore_color.rgb  = RGBColor(250, 200, 200)
                
def replace_text_in_shape(shape, nc_result: dict):
    """Replace tag text in a shape"""
    for match, replacement in nc_result.items():
        if shape.has_text_frame:
            if (shape.text.find(match)) != -1:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    whole_text = "".join(run.text for run in paragraph.runs)
                    whole_text = whole_text.replace(str(match), str(replacement))
                    for idx, run in enumerate(paragraph.runs):
                        if idx != 0:
                            p = paragraph._p
                            p.remove(run._r)
                    if bool(paragraph.runs):
                        font_size = paragraph.runs[0].font.size
                        paragraph.runs[0].text = whole_text
                        paragraph.runs[0].font.color.rgb = RGBColor(0x0, 0x0, 0x0)
                        paragraph.runs[0].font.size = font_size


def replace_images_in_slide(slide, nc_result: dict):
    """Replace images in slide"""
    for shape in slide.shapes:
        for key in nc_result:
            if 'Image' in key or  'Score' in key:
                image_name = key.replace('{','').replace('}','')
                image_path = nc_result[key]
                if len(image_path) > 4 and image_name in shape.name:
                    slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
                    slide.shapes.element.remove(shape.element)   
       

def delete_prev_nc_images(nc_result):
    for key in nc_result:
        if 'Image' in key or 'Score' in key:
            image_path = nc_result[key]
            try:
                os.remove(image_path)
            except Exception as e:
                pass

def get_column_number_from_range(cell_range, search_value):
    for cell in cell_range:
        cell_value = "".join(re.findall("[a-zA-Z0-9-+]+", str(cell.Value)))
        search_value = "".join(re.findall("[a-zA-Z0-9-+]+", str(search_value)))
        if cell_value.upper() == search_value.upper():
            return cell.Column
    return 0
    
def is_float(string):
    try:
        float(string)
        return True
    except ValueError:
        return False

def check_for_files(folder_path, pattern1, pattern2=None, pattern3=None):
    """Search and list pdf files in the given folder"""
    list1 = glob.glob(os.path.join(folder_path, pattern1)) if pattern1 else []
    list2 = glob.glob(os.path.join(folder_path, pattern2)) if pattern2 else []
    list3 = glob.glob(os.path.join(folder_path, pattern3)) if pattern3 else []
    return list1 + list2 + list3


def delete_temp_files():
    try:
        temp_folder = os.environ['TEMP']
        for filename in os.listdir(temp_folder):
            file_path = os.path.join(temp_folder, filename)
            try:
                if os.path.isfile(file_path) or os.path.islink(file_path):
                    os.unlink(file_path)
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
            except:
                pass
    except:
        pass


def save_shape_as_image(shape, img_name):
    loop_count = 0
    while True:
        loop_count += 1
        try:
            shape.Copy()
            image = ImageGrab.grabclipboard()
            if image.mode in ("RGBA", "P"):
                image = image.convert("RGB")
            image.save(img_name, 'png')
            return True
        except Exception as e:
            print("Unable to save object '{}' as image. Retry Loop {}".format(img_name, loop_count))
            clear_clip_board()
            sleep(1)
            if loop_count > 5:
                break
    return False


def save_range_as_image(copy_range, img_name):
    loop_count = 0
    while True:
        loop_count += 1
        try:
            
            copy_range.Copy()
            image = ImageGrab.grabclipboard()
            if image.mode in ("RGBA", "P"):
                image = image.convert("RGB")
            image.save(img_name, 'png')
            return True
        except Exception as e:
            print("Unable to save object '{}' as image. Retry Loop {}".format(img_name, loop_count))
            clear_clip_board()
            sleep(1)
            if loop_count > 5:
                break
    return False


def isDigit(x):
    try:
        float(x)
        return True
    except:
        return False


def clear_clip_board():
    loop_count = 0
    while True:
        loop_count += 1
        try:
            ctypes.windll.user32.OpenClipboard(None)
            ctypes.windll.user32.EmptyClipboard()
            ctypes.windll.user32.CloseClipboard()
            break
        except:
            sleep(5)
            if loop_count > 5:
                break
