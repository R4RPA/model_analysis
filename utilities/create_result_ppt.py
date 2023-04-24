import datetime
import utilities.support_functions as sup_functions
from pptx import Presentation
from pptx.table import Table
import os.path
import math


def create_result_ppt(template_path, out_path, result_dict):
    """Create results report using the template given"""

    """Load the report template"""
    prs = Presentation(template_path)
    slides = [slide for slide in prs.slides]
    tip_fix_results_list = result_dict['tip_fix_list']
    tip_free_results_list = result_dict['tip_free_list']
    results = [tip_fix_results_list, tip_free_results_list]
    shape_slide_list = [8, 13]
    table_name_list = ['Table 1', 'Table 2']
    rst_type_list = ['tip_fix', 'tip_free']
    header_shape_list = ['Tip_Fix_Slide_Header', 'Tip_Free_Slide_Header']
    txt_match_list = ['{T1_PageNum}', '{T2_PageNum}']
    slides_added = 0
    atlest_one_result = False
    slides_to_delete = []
    for index, results_list in enumerate(results):
        if len(results_list) > 0:
            atlest_one_result = True
            number_of_modes = len(results_list)
            number_of_slides = math.ceil(number_of_modes / 6)

            """Update Frequncy List"""
            summary_table = sup_functions.get_object_by_name(slides[shape_slide_list[index]-1], table_name_list[index])
            summary_table = sup_functions.add_multiple_rows(table=summary_table, number_of_rows=number_of_modes - 1)
            update_summary_table(summary_table, results_list)

            block_layout = prs.slide_masters[3].slide_layouts[0]

            copy_slide_index = shape_slide_list[index] + slides_added
            dest_slide_index = copy_slide_index
            slides_to_delete.append(copy_slide_index)
            start_index = 0
            for slide_num in range(1, number_of_slides+1):
                slice_ = results_list[start_index:start_index + 6]
                start_index += 6
                dest_slide_index += 1
                prs = sup_functions.duplicate_slide(prs, block_layout, copy_slide_index, dest_slide_index, rst_type_list[index], slice_)

                header_shape = sup_functions.get_object_by_name(prs.slides[dest_slide_index], header_shape_list[index])
                match_dict = {txt_match_list[index]: slide_num}
                sup_functions.replace_text_in_shape(header_shape, match_dict)

            slides_added = number_of_slides

    if atlest_one_result:
        """Delete template mode slides"""
        prs = sup_functions.delete_slide(prs, slides_to_delete)

        """Save the report ppt"""
        timestamp = datetime.datetime.now().strftime("%Y%m%d%H%M%S")
        output_file_path = os.path.join(out_path, f'rst_result_{timestamp}.pptx')
        prs.save(output_file_path)


def update_summary_table(table: Table, replacements: list):
    """Update Summary table"""
    for row_num, nc_result in enumerate(replacements):
        table = sup_functions.update_table_row(table, nc_result, row_num, font_size=8)
    return table


def main():
    root_dir = os.path.dirname(os.path.abspath(__file__))
    root_dir = os.path.abspath(os.path.join(root_dir, os.pardir))
    template_path = os.path.join(root_dir, 'data/template/modal_analysis_final_report_template.pptx')
    output_path = os.path.join(root_dir, 'data/output_reports')
    if not os.path.exists(output_path):
        os.makedirs(output_path)

    results_list = [
                          {
                            "mode": 1,
                            "frequency": 71.181,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 2,
                            "frequency": 71.182,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 3,
                            "frequency": 71.183,
                            "image": "data/result_images/tip_fix_mode_fig_2.png"
                          },
                          {
                            "mode": 4,
                            "frequency": 71.184,
                            "image": "data/result_images/tip_fix_mode_fig_3.png"
                          },
                          {
                            "mode": 5,
                            "frequency": 71.185,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 6,
                            "frequency": 71.186,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 7,
                            "frequency": 71.187,
                            "image": "data/result_images/tip_fix_mode_fig_2.png"
                          },
                          {
                            "mode": 8,
                            "frequency": 71.188,
                            "image": "data/result_images/tip_fix_mode_fig_3.png"
                          },
                          {
                            "mode": 9,
                            "frequency": 71.189,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 10,
                            "frequency": 71.190,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 2,
                            "frequency": 71.182,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 3,
                            "frequency": 71.183,
                            "image": "data/result_images/tip_fix_mode_fig_2.png"
                          },
                          {
                            "mode": 4,
                            "frequency": 71.184,
                            "image": "data/result_images/tip_fix_mode_fig_3.png"
                          },
                          {
                            "mode": 5,
                            "frequency": 71.185,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 6,
                            "frequency": 71.186,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 7,
                            "frequency": 71.187,
                            "image": "data/result_images/tip_fix_mode_fig_2.png"
                          },
                          {
                            "mode": 8,
                            "frequency": 71.188,
                            "image": "data/result_images/tip_fix_mode_fig_3.png"
                          },
                          {
                            "mode": 9,
                            "frequency": 71.189,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 10,
                            "frequency": 71.190,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 2,
                            "frequency": 71.182,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 3,
                            "frequency": 71.183,
                            "image": "data/result_images/tip_fix_mode_fig_2.png"
                          },
                          {
                            "mode": 4,
                            "frequency": 71.184,
                            "image": "data/result_images/tip_fix_mode_fig_3.png"
                          },
                          {
                            "mode": 5,
                            "frequency": 71.185,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 6,
                            "frequency": 71.186,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 7,
                            "frequency": 71.187,
                            "image": "data/result_images/tip_fix_mode_fig_2.png"
                          },
                          {
                            "mode": 8,
                            "frequency": 71.188,
                            "image": "data/result_images/tip_fix_mode_fig_3.png"
                          },
                          {
                            "mode": 9,
                            "frequency": 71.189,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 10,
                            "frequency": 71.190,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 2,
                            "frequency": 71.182,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 3,
                            "frequency": 71.183,
                            "image": "data/result_images/tip_fix_mode_fig_2.png"
                          },
                          {
                            "mode": 4,
                            "frequency": 71.184,
                            "image": "data/result_images/tip_fix_mode_fig_3.png"
                          },
                          {
                            "mode": 5,
                            "frequency": 71.185,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 6,
                            "frequency": 71.186,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          },
                          {
                            "mode": 7,
                            "frequency": 71.187,
                            "image": "data/result_images/tip_fix_mode_fig_2.png"
                          },
                          {
                            "mode": 8,
                            "frequency": 71.188,
                            "image": "data/result_images/tip_fix_mode_fig_3.png"
                          },
                          {
                            "mode": 9,
                            "frequency": 71.189,
                            "image": "data/result_images/tip_fix_mode_fig_0.png"
                          },
                          {
                            "mode": 10,
                            "frequency": 71.190,
                            "image": "data/result_images/tip_fix_mode_fig_1.png"
                          }
                        ]
    result_dict = {'tip_fix_list': results_list, 'tip_free_list': results_list}
    create_result_ppt(template_path, output_path, result_dict)
    print('Check : ', output_path, 'for reports')


if __name__ == '__main__':
    main()

